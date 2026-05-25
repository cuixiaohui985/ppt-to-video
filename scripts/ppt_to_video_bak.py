#!/usr/bin/env python3
"""
PPT to Video Converter（优化版 v0.2）
将 PPT 文件转换为 MP4 视频，支持：
1. 智能时长——封面/过渡页 1.5s，正文页按 180字/分计算（3字/秒），范围 6~8s
2. 翻页动画——支持向上翻页（slideup）和右下到左上翻页（diagtl）
3. 打字机效果——正文内容页文字逐字渐入（通过 ffmpeg drawtext 实现）

依赖：
  pip install python-pptx pillow pywin32 pdf2image
  系统需安装 ffmpeg（用于视频合成）、ffprobe（获取音频时长）
  可选：LibreOffice（跨平台渲染后备方案）

用法：
  python ppt_to_video.py <ppt_path> <audio_path> [选项]
"""

import argparse
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("缺少 python-pptx 库，请运行：pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print("缺少 pillow 库，请运行：pip install pillow")
    sys.exit(1)


# ────────────────────────────────────────────
# 1. 从 PPT 提取文字内容，计算每页时长
# ────────────────────────────────────────────

def extract_slide_text(ppt_path: str) -> list:
    """
    返回 [(slide_idx, text, char_count), ...]
    """
    prs = Presentation(ppt_path)
    results = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        full_text = " ".join(texts)
        char_count = len(re.sub(r"\s+", "", full_text))
        results.append((idx, full_text, char_count))
    return results


def classify_slide(idx: int, total: int, char_count: int) -> str:
    """
    判断页面类型：
  - 'cover'  ：封面/尾页（首页、末页）或字数 < 50
  - 'content'：正文内容页
    """
    if idx == 1 or idx == total or char_count < 50:
        return "cover"
    return "content"


def calc_durations(
    slide_texts: list,
    cover_dur: float = 1.5,
    back_cover_dur: float = 4.0,
    content_cps: float = 3.0,   # 180字/分 = 3字/秒
    content_min: float = 6.0,
    content_max: float = 8.0,
) -> list:
    """
    智能计算每页时长。
    封面/过渡页：cover_dur（默认 1.5s），末页可用 back_cover_dur 单独控制
    正文页：clamp(char_count / content_cps, content_min, content_max)
    返回 [(slide_idx, duration_sec, slide_type), ...]
    """
    total = len(slide_texts)
    durations = []
    for idx, _, char_count in slide_texts:
        stype = classify_slide(idx, total, char_count)
        if stype == "cover":
            if idx == total:
                dur = back_cover_dur   # 封底单独时长（默认 4s）
            else:
                dur = cover_dur
        else:
            dur = char_count / content_cps
            dur = max(content_min, min(content_max, dur))
        # 取整到 0.5 秒，便于 ffmpeg 处理
        dur = round(dur * 2) / 2
        durations.append((idx, dur, stype))
    return durations


# ────────────────────────────────────────────
# 2. 将 PPT 每页渲染为图片（优化渲染质量）
# ────────────────────────────────────────────

def _set_ppt_export_resolution():
    """设置 PowerPoint 导出分辨率为 300DPI（仅 Windows）。
    修改 HKCU\\Software\\Microsoft\\Office\\16.0\\PowerPoint\\Options
    ExportBitmapResolution = 300
    """
    try:
        import winreg
        # 尝试常见 Office 版本
        for ver in ["16.0", "15.0", "14.0", "12.0"]:
            key_path = rf"Software\Microsoft\Office\{ver}\PowerPoint\Options"
            try:
                key = winreg.CreateKeyEx(
                    winreg.HKEY_CURRENT_USER, key_path, 0,
                    winreg.KEY_SET_VALUE | winreg.KEY_QUERY_VALUE
                )
                # 读取当前值
                try:
                    current, _ = winreg.QueryValueEx(key, "ExportBitmapResolution")
                except FileNotFoundError:
                    current = 0

                if current != 300:
                    winreg.SetValueEx(key, "ExportBitmapResolution", 0, winreg.REG_DWORD, 300)
                    print("  ✓ 已设置 PowerPoint 导出分辨率为 300DPI")

                winreg.CloseKey(key)
                return  # 成功设置，退出
            except Exception:
                continue
    except Exception as e:
        print(f"  ⚠ 设置导出分辨率失败（不影响功能）：{e}")


def convert_ppt_to_images_via_com(ppt_path: str, out_dir: str) -> list:
    """使用 PowerPoint COM 自动化（仅 Windows），已设置 300DPI 导出。"""
    try:
        import win32com.client
    except ImportError:
        raise ImportError("需要 pywin32。请运行：pip install pywin32")

    _set_ppt_export_resolution()

    ppt_path = os.path.abspath(ppt_path)
    out_dir = os.path.abspath(out_dir)
    os.makedirs(out_dir, exist_ok=True)

    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    try:
        powerpoint.Visible = False
    except Exception:
        pass  # PowerPoint 2007 不支持此属性，忽略
    try:
        pres = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
        image_paths = []
        for i in range(1, pres.Slides.Count + 1):
            path = os.path.join(out_dir, f"slide_{i:03d}.png")
            try:
                # PowerPoint 2010+ 支持宽高参数
                pres.Slides(i).Export(path, "PNG", 1920, 1080)
            except Exception:
                # PowerPoint 2007 仅支持 FileName + FilterName
                pres.Slides(i).Export(path, "PNG")
            image_paths.append(path)
        pres.Close()
    finally:
        powerpoint.Quit()
    return image_paths


def convert_ppt_to_images_via_libreoffice(ppt_path: str, out_dir: str) -> list:
    """
    使用 LibreOffice 转换（跨平台，高质量）。
    方案A（优先）：PDF 导出 → pdf2image 转 PNG（300DPI，质量最高）
    方案B（后备）：直接 PNG 导出
    """
    ppt_path = os.path.abspath(ppt_path)
    out_dir = os.path.abspath(out_dir)
    os.makedirs(out_dir, exist_ok=True)
    errors = []

    # ── 方案A：PDF 导出 → pdf2image（质量最高）───
    pdf_dir = os.path.join(out_dir, "pdf")
    os.makedirs(pdf_dir, exist_ok=True)
    cmd_pdf = [
        "soffice", "--headless", "--invisible",
        "--convert-to", "pdf:writer_pdf_Export",
        "--outdir", pdf_dir, ppt_path
    ]
    result_a = subprocess.run(cmd_pdf, capture_output=True, text=True)
    pdf_files = sorted([f for f in os.listdir(pdf_dir) if f.endswith(".pdf")])

    if result_a.returncode == 0 and pdf_files:
        pdf_path = os.path.join(pdf_dir, pdf_files[0])
        try:
            from pdf2image import convert_from_path
            # DPI=300 保证最高质量，后续 ffmpeg 会缩放到 1920x1080
            images = convert_from_path(pdf_path, dpi=300)
            image_paths = []
            for i, img in enumerate(images):
                path = os.path.join(out_dir, f"slide_{i+1:03d}.png")
                img.save(path, "PNG")
                image_paths.append(path)
            print("  ✓ LibreOffice PDF→PNG（300DPI）")
            return image_paths
        except Exception as e:
            errors.append(f"PDF→PNG 失败：{e}")

    # ── 方案B：直接 PNG 导出（后备）───
    png_dir = os.path.join(out_dir, "png")
    os.makedirs(png_dir, exist_ok=True)
    cmd_png = [
        "soffice", "--headless", "--invisible",
        "--convert-to", "png",
        "--outdir", png_dir, ppt_path
    ]
    result_b = subprocess.run(cmd_png, capture_output=True, text=True)
    png_files = sorted([f for f in os.listdir(png_dir) if f.endswith(".png")])

    if result_b.returncode == 0 and png_files:
        image_paths = []
        for f in png_files:
            src = os.path.join(png_dir, f)
            dst = os.path.join(out_dir, f)
            shutil.move(src, dst)
            image_paths.append(dst)
        print("  ✓ LibreOffice 直接 PNG 导出")
        return image_paths

    # 两种方案都失败
    if result_a.returncode != 0:
        errors.append(f"PDF 导出失败：{result_a.stderr[-300:]}")
    if result_b.returncode != 0:
        errors.append(f"PNG 导出失败：{result_b.stderr[-300:]}")
    raise RuntimeError(
        "LibreOffice 转换失败：\n" + "\n".join(errors) +
        "\n请确保已安装 LibreOffice 并将 soffice 加入 PATH。"
    )


def convert_ppt_to_images(ppt_path: str, out_dir: str) -> list:
    """
    自动选择可用方案，优先级：
    1. COM（Windows，质量最高，已设 300DPI）
    2. LibreOffice PDF→PNG（跨平台，质量高）
    3. LibreOffice 直接 PNG（后备）
    """
    image_paths = []
    errors = []

    # 方案1：COM（仅 Windows）
    if sys.platform == "win32":
        try:
            image_paths = convert_ppt_to_images_via_com(ppt_path, out_dir)
            if image_paths:
                print(f"  ✓ COM 渲染（PowerPoint 原生，300DPI）")
                return image_paths
        except Exception as e:
            errors.append(f"COM 方案失败：{e}")

    # 方案2：LibreOffice（PDF→PNG 优先，然后直接 PNG）
    try:
        image_paths = convert_ppt_to_images_via_libreoffice(ppt_path, out_dir)
        if image_paths:
            return image_paths
    except Exception as e:
        errors.append(f"LibreOffice 方案失败：{e}")

    raise RuntimeError("所有 PPT 转图片方案均失败：\n" + "\n".join(errors))


def calc_total_duration(durations: list, transition_dur: float) -> float:
    """计算 xfade 合成后的正确总时长。"""
    video_dur = sum(d for _, d, _ in durations)
    n = len(durations)
    if n <= 1 or transition_dur <= 0:
        return video_dur
    return round(video_dur - (n - 1) * transition_dur, 2)


# ────────────────────────────────────────────
# 3. 翻页动画（ffmpeg xfade）
# ────────────────────────────────────────────
# xfade transition 名称映射：
#   "slideup"   → 向上翻页
#   "diagtl"     → 从右下向左上翻页（diagonal ↗）
#   "fade"       → 淡入淡出
#   "none"       → 无动画

XFADE_NAMES = {
    "slideup": "slideup",
    "up": "slideup",
    "diagtl": "diagtl",
    "diagonal": "diagtl",
    "fade": "fade",
    "none": None,
}


def build_video_with_transitions(
    image_paths: list,
    durations: list,
    audio_path: str,
    output_path: str,
    transition: str = "slideup",
    transition_dur: float = 0.5,
    fps: int = 25,
) -> None:
    """
    使用 ffmpeg xfade 滤镜合成带翻页动画的视频。

    正确公式（经过实测验证）：
      offset_i = sum(durations[0..i]) - (i+1) * transition_dur
      总时长 = sum(all durations) - (n-1) * transition_dur
    """
    if transition == "none" or transition is None:
        build_video_simple(image_paths, durations, audio_path, output_path, fps)
        return

    temp_dir = os.path.dirname(output_path)
    clip_dir = os.path.join(temp_dir, "clips")
    os.makedirs(clip_dir, exist_ok=True)

    # ── Step A：生成每段图片对应的视频片段 ──
    clip_paths = []
    for i, (img_path, (_, dur, _)) in enumerate(zip(image_paths, durations)):
        clip_path = os.path.join(clip_dir, f"clip_{i:03d}.mp4")
        cmd = [
            "ffmpeg", "-y",
            "-loop", "1",
            "-i", img_path,
            "-c:v", "libx264", "-preset", "fast", "-crf", "23",
            "-t", str(dur),
            "-vf", f"fps={fps},scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2",
            "-pix_fmt", "yuv420p",
            "-an",
            clip_path
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f"生成片段失败（第 {i+1} 页）：{result.stderr}")
        clip_paths.append(clip_path)

    # ── Step B：计算 xfade offset（正确公式）───
    # offset_i = sum(durations[0..i]) - (i+1) * transition_dur
    offsets = []
    cumulative = 0.0
    for i, (_, dur, _) in enumerate(durations[:-1]):
        cumulative += dur
        offset = round(cumulative - (i + 1) * transition_dur, 2)
        offsets.append(offset)

    if len(clip_paths) == 1:
        intermediate = clip_paths[0]
    else:
        # 构建 xfade filter_complex 字符串
        filter_parts = []
        for i in range(1, len(clip_paths)):
            offset = offsets[i - 1]
            if i == 1:
                filter_parts.append(
                    f"[0:v][1:v]xfade=transition={transition}:duration={transition_dur}:offset={offset}[v1]"
                )
            else:
                prev = f"[v{i-1}]"
                filter_parts.append(
                    f"{prev}[{i}:v]xfade=transition={transition}:duration={transition_dur}:offset={offset}[v{i}]"
                )

        filter_complex = "; ".join(filter_parts)
        prev_label = f"[v{len(clip_paths)-1}]"

        intermediate_path = os.path.join(clip_dir, "intermediate.mp4")
        cmd = ["ffmpeg", "-y"]
        for cp in clip_paths:
            cmd += ["-i", cp]
        cmd += [
            "-filter_complex", filter_complex,
            "-map", prev_label,
            "-c:v", "libx264", "-preset", "fast", "-crf", "23",
            "-pix_fmt", "yuv420p",
            "-an",
            intermediate_path
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f"xfade 合成失败：{result.stderr}")
        intermediate = intermediate_path

    # ── Step C：叠加背景音乐（循环）───
    # 正确做法：用 calc_total_duration() 计算总时长，传给 -t
    audio_dur = get_audio_duration(audio_path)
    correct_total = calc_total_duration(durations, transition_dur)
    loop_count = max(1, math.ceil(correct_total / audio_dur)) if audio_dur > 0 else 1

    cmd = [
        "ffmpeg", "-y",
        "-i", intermediate,
        "-stream_loop", str(loop_count),
        "-i", audio_path,
        "-map", "0:v",
        "-map", "1:a",
        "-c:v", "copy",
        "-c:a", "aac", "-b:a", "128k",
        "-t", str(correct_total),
        "-shortest",
        output_path
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"叠加音频失败：{result.stderr}")

    # 清理片段文件
    for cp in clip_paths:
        try: os.remove(cp)
        except: pass
    try:
        if 'intermediate' in locals() and os.path.exists(intermediate):
            os.remove(intermediate)
    except: pass


def build_video_simple(
    image_paths: list,
    durations: list,
    audio_path: str,
    output_path: str,
    fps: int = 25,
) -> None:
    """无动画的退路方案（concat demuxer）。"""
    temp_dir = os.path.dirname(output_path)
    concat_path = os.path.join(temp_dir, "concat.txt")
    with open(concat_path, "w", encoding="utf-8") as f:
        for (_, dur, _), img_path in zip(durations, image_paths):
            f.write(f"file '{img_path}'\n")
            f.write(f"duration {dur}\n")
        if image_paths:
            f.write(f"file '{image_paths[-1]}'\n")

    video_dur = sum(d for _, d, _ in durations)
    audio_dur = get_audio_duration(audio_path)
    loop_count = math.ceil(video_dur / audio_dur) if audio_dur > 0 else 1

    cmd = [
        "ffmpeg", "-y",
        "-f", "concat", "-safe", "0", "-i", concat_path,
        "-stream_loop", str(loop_count),
        "-i", audio_path,
        "-map", "0:v", "-map", "1:a",
        "-c:v", "libx264", "-preset", "medium", "-crf", "23",
        "-vf", f"fps={fps},scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2",
        "-t", str(video_dur + 1),
        "-c:a", "aac", "-b:a", "128k",
        "-shortest",
        output_path
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        alt_cmd = [
            "ffmpeg", "-y",
            "-f", "concat", "-safe", "0", "-i", concat_path,
            "-stream_loop", str(loop_count),
            "-i", audio_path,
            "-map", "0:v", "-map", "1:a",
            "-c:v", "libx264", "-preset", "fast", "-crf", "23",
            "-vf", f"fps={fps}",
            "-c:a", "aac",
            "-t", str(video_dur + 0.5),
            "-shortest",
            output_path
        ]
        result2 = subprocess.run(alt_cmd, capture_output=True, text=True)
        if result2.returncode != 0:
            raise RuntimeError(
                f"ffmpeg 执行失败：\n{result.stderr}\n\n{result2.stderr}"
            )


# ────────────────────────────────────────────
# 4. 打字机效果（ffmpeg drawtext）
# ────────────────────────────────────────────

def get_audio_duration(audio_path: str) -> float:
    """使用 ffprobe 获取音频时长（秒）。"""
    cmd = [
        "ffprobe", "-v", "error",
        "-show_entries", "format=duration",
        "-of", "default=noprint_wrappers=1:nokey=1",
        audio_path
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    return float(result.stdout.strip())


def detect_system_font() -> str:
    """自动检测系统可用的中文字体路径。"""
    candidates = [
        # Windows
        "C:/Windows/Fonts/msyh.ttc",       # 微软雅黑
        "C:/Windows/Fonts/simsun.ttc",      # 宋体
        "C:/Windows/Fonts/simhei.ttf",      # 黑体
        # Linux
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        # macOS
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/System/Library/Fonts/PingFang.ttc",
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return "arial.ttf"   # 最后退路


def escape_text_for_ffmpeg(text: str) -> str:
    """转义 ffmpeg drawtext 中的特殊字符。"""
    text = text.replace("\\", "\\\\")
    text = text.replace("'", "'\\''")
    text = text.replace(":", "\\:")
    text = text.replace(",", "\\,")
    return text


def build_video_with_typewriter(
    image_paths: list,
    durations: list,
    slide_texts: list,
    audio_path: str,
    output_path: str,
    transition: str = "slideup",
    transition_dur: float = 0.5,
    fps: int = 25,
) -> None:
    """
    带打字机效果（正文页）的视频合成。
    封面页：直接使用原图（无效果）
    正文页：用 ffmpeg drawtext 实现逐字渐入
    """
    temp_dir = os.path.dirname(output_path)
    work_dir = os.path.join(temp_dir, "typewriter_work")
    os.makedirs(work_dir, exist_ok=True)

    font_path = detect_system_font()
    print(f"  使用字体：{font_path}")

    # ── 为每一页生成对应的视频片段 ──
    clip_paths = []
    for i, (img_path, (idx, dur, stype)) in enumerate(zip(image_paths, durations)):
        clip_path = os.path.join(work_dir, f"clip_{i:03d}.mp4")
        text = slide_texts[idx - 1][1] if idx <= len(slide_texts) else ""

        if stype == "cover" or not text.strip():
            # 封面/无文字 → 直接生成静态片段
            cmd = [
                "ffmpeg", "-y",
                "-loop", "1", "-i", img_path,
                "-c:v", "libx264", "-preset", "fast", "-crf", "23",
                "-t", str(dur),
                "-vf", f"fps={fps},scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2",
                "-pix_fmt", "yuv420p",
                "-an",
                clip_path
            ]
        else:
            # 正文页 → 打字机效果
            safe_text = escape_text_for_ffmpeg(text[:200])  # 限制长度
            chars_to_show = len(safe_text)
            char_rate = chars_to_show / (dur * fps) if dur > 0 else 1
            draw_expr = f"min(int(n*{char_rate:.4f}),{chars_to_show})"

            vf = (
                f"fps={fps},"
                f"scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2,"
                f"drawtext=fontfile='{font_path}':text='{safe_text}':"
                f"fontsize=48:fontcolor=white:x=(w-text_w)/2:y=(h-text_h)/2:"
                f"box=1:boxcolor=black@0.5:boxborderw=5:"
                f"enable='gte(t,0)':draw='{draw_expr}'"
            )

            cmd = [
                "ffmpeg", "-y",
                "-loop", "1", "-i", img_path,
                "-c:v", "libx264", "-preset", "fast", "-crf", "23",
                "-t", str(dur),
                "-vf", vf,
                "-pix_fmt", "yuv420p",
                "-an",
                clip_path
            ]

        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            print(f"  ⚠ 第 {idx} 页渲染失败，使用普通片段")
            fallback_cmd = [
                "ffmpeg", "-y",
                "-loop", "1", "-i", img_path,
                "-c:v", "libx264", "-preset", "fast", "-crf", "23",
                "-t", str(dur),
                "-vf", f"fps={fps},scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2",
                "-pix_fmt", "yuv420p",
                "-an",
                clip_path
            ]
            subprocess.run(fallback_cmd, check=True, capture_output=True)

        clip_paths.append(clip_path)

    # ── 串联所有片段（xfade）并叠加音频 ──
    if transition == "none" or len(clip_paths) == 1:
        intermediate = clip_paths[0] if clip_paths else None
    else:
        offsets = []
        cumulative = 0.0
        for i, (_, dur, _) in enumerate(durations[:-1]):
            cumulative += dur
            offsets.append(round(cumulative - (i + 1) * transition_dur, 2))

        filter_parts = []
        for i in range(1, len(clip_paths)):
            offset = offsets[i - 1]
            if i == 1:
                filter_parts.append(
                    f"[0:v][1:v]xfade=transition={transition}:duration={transition_dur}:offset={offset}[v1]"
                )
            else:
                prev = f"[v{i-1}]"
                filter_parts.append(
                    f"{prev}[{i}:v]xfade=transition={transition}:duration={transition_dur}:offset={offset}[v{i}]"
                )

        filter_complex = "; ".join(filter_parts)
        prev_label = f"[v{len(clip_paths)-1}]" if len(clip_paths) > 1 else "[0:v]"

        intermediate_path = os.path.join(work_dir, "intermediate.mp4")
        cmd = ["ffmpeg", "-y"]
        for cp in clip_paths:
            cmd += ["-i", cp]
        cmd += [
            "-filter_complex", filter_complex,
            "-map", prev_label,
            "-c:v", "libx264", "-preset", "fast", "-crf", "23",
            "-pix_fmt", "yuv420p",
            "-an",
            intermediate_path
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f"xfade 合成失败：{result.stderr}")
        intermediate = intermediate_path

    # 叠加音频
    if intermediate:
        correct_total = calc_total_duration(durations, transition_dur)
        audio_dur = get_audio_duration(audio_path)
        loop_count = max(1, math.ceil(correct_total / audio_dur)) if audio_dur > 0 else 1

        cmd = [
            "ffmpeg", "-y",
            "-i", intermediate,
            "-stream_loop", str(loop_count),
            "-i", audio_path,
            "-map", "0:v", "-map", "1:a",
            "-c:v", "copy",
            "-c:a", "aac", "-b:a", "128k",
            "-t", str(correct_total),
            "-shortest",
            output_path
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f"叠加音频失败：{result.stderr}")

    # 清理
    for cp in clip_paths:
        try: os.remove(cp)
        except: pass


# ────────────────────────────────────────────
# 5. 主流程
# ────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="PPT 转视频工具（优化版 v0.2）")
    parser.add_argument("ppt", help="PPT/PPTX 文件路径")
    parser.add_argument("audio", help="背景音乐文件路径（MP3/WAV/M4A 等）")
    parser.add_argument("--output", default=None, help="输出 MP4 路径（默认同目录）")
    parser.add_argument("--transition", default="slideup",
                        choices=["slideup", "diagtl", "fade", "none"],
                        help="翻页动画类型（默认 slideup）")
    parser.add_argument("--transition-dur", type=float, default=0.5,
                        help="翻页动画时长（秒），默认 0.5")
    parser.add_argument("--typewriter", action="store_true",
                        help="为正文页启用打字机文字渐入效果")
    parser.add_argument("--cover-dur", type=float, default=1.5,
                        help="封面/过渡页时长（秒），默认 1.5")
    parser.add_argument("--back-cover-dur", type=float, default=4.0,
                        help="尾页（封底）时长（秒），默认 4.0")
    parser.add_argument("--content-cps", type=float, default=3.0,
                        help="正文页语速（字/秒），默认 3（=180字/分）")
    parser.add_argument("--content-min", type=float, default=6.0,
                        help="正文页最短时长（秒），默认 6")
    parser.add_argument("--content-max", type=float, default=8.0,
                        help="正文页最长时长（秒），默认 8")
    parser.add_argument("--fps", type=int, default=25, help="输出视频帧率，默认 25")
    parser.add_argument("--keep-temp", action="store_true", help="保留临时文件（调试用）")
    args = parser.parse_args()

    ppt_path = os.path.abspath(args.ppt)
    audio_path = os.path.abspath(args.audio)

    if not os.path.exists(ppt_path):
        print(f"错误：PPT 文件不存在：{ppt_path}")
        sys.exit(1)
    if not os.path.exists(audio_path):
        print(f"错误：音频文件不存在：{audio_path}")
        sys.exit(1)

    if args.output:
        output_path = os.path.abspath(args.output)
    else:
        base = os.path.splitext(ppt_path)[0]
        output_path = base + ".mp4"

    temp_dir = tempfile.mkdtemp(prefix="ppt2video_")
    print(f"临时目录：{temp_dir}")

    try:
        # Step 1: 提取文字
        print("Step 1/5：提取 PPT 文字内容...")
        slide_texts = extract_slide_text(ppt_path)
        durations = calc_durations(
            slide_texts,
            cover_dur=args.cover_dur,
            back_cover_dur=args.back_cover_dur,
            content_cps=args.content_cps,
            content_min=args.content_min,
            content_max=args.content_max,
        )
        total_dur = sum(d for _, d, _ in durations)
        print(f"  共 {len(durations)} 页，总时长：{total_dur:.1f} 秒")
        for idx, dur, stype in durations:
            chars = slide_texts[idx - 1][2]
            tag = "（封面）" if stype == "cover" else "（正文）"
            print(f"  第 {idx} 页：{chars} 字 → {dur} 秒 {tag}")

        # Step 2: PPT -> 图片
        print("Step 2/5：将 PPT 每页转为图片...")
        image_dir = os.path.join(temp_dir, "images")
        os.makedirs(image_dir, exist_ok=True)
        image_paths = convert_ppt_to_images(ppt_path, image_dir)
        print(f"  已生成 {len(image_paths)} 张图片")

        # Step 3: 获取音频时长
        print("Step 3/5：分析背景音乐...")
        audio_dur = get_audio_duration(audio_path)
        correct_total = calc_total_duration(durations, args.transition_dur)
        print(f"  音频时长：{audio_dur:.1f} 秒，需要循环 {math.ceil(correct_total/audio_dur)} 次")

        # Step 4: 合成视频
        transition_name = XFADE_NAMES.get(args.transition, "slideup")
        if args.typewriter:
            print("Step 4/5：合成视频（打字机效果 + 翻页动画）...")
            build_video_with_typewriter(
                image_paths, durations, slide_texts,
                audio_path, output_path,
                transition=transition_name,
                transition_dur=args.transition_dur,
                fps=args.fps,
            )
        elif transition_name and transition_name != "none":
            print(f"Step 4/5：合成视频（翻页动画：{transition_name}）...")
            build_video_with_transitions(
                image_paths, durations, audio_path, output_path,
                transition=transition_name,
                transition_dur=args.transition_dur,
                fps=args.fps,
            )
        else:
            print("Step 4/5：合成视频（无翻页动画）...")
            build_video_simple(
                image_paths, durations, audio_path, output_path, args.fps
            )

        # Step 5: 完成
        print(f"\n✅ 视频已生成：{output_path}")
        print(f"   分辨率：1920×1080，帧率：{args.fps}fps")
        print(f"   总时长：{total_dur:.1f} 秒")
        if args.typewriter:
            print(f"   打字机效果：已启用（正文页）")
        if transition_name and transition_name != "none":
            print(f"   翻页动画：{transition_name}")
        print(f"   渲染方案：{_get_render_method()}")

    finally:
        if not args.keep_temp:
            shutil.rmtree(temp_dir, ignore_errors=True)
        else:
            print(f"临时文件保留在：{temp_dir}")


def _get_render_method() -> str:
    """返回当前使用的渲染方案（用于日志）。"""
    if sys.platform == "win32":
        return "COM (PowerPoint 原生)"
    return "LibreOffice PDF→PNG (300DPI)"


if __name__ == "__main__":
    main()
